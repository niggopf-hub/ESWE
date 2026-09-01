#!/usr/bin/env python3
"""Inventar einer Unterlage: Folien, Layouts, Platzhalter, Textboxen, Diagramme.

Damit findest du die Box, in die geschrieben werden soll — und siehst, welche
Pflicht-Platzhalter noch leer sind.

    python3 inspect_deck.py deck.pptx              Uebersicht ueber alle Folien
    python3 inspect_deck.py deck.pptx --slide 12   eine Folie im Detail
    python3 inspect_deck.py deck.pptx --check      Pruefmodus vor Uebergabe

Benoetigt: pip install python-pptx
"""
import argparse
import re
import sys

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("Fehlt: pip install python-pptx")

# Platzhalter-Indizes der Inhalt_0-Familie (siehe references/powerpoint.md)
ROLLEN = {
    0: "Titel (Aussagesatz)",
    10: "Kapitelnummer",
    11: "Kapiteltitel",
    12: "Bereichsueberschrift",
    13: "Kapitelkolumne",
    14: "Quellenzeile",
    15: "Fussnotenzeile",
}
PFLICHT_INHALT = {0: "Titel", 13: "Kapitelkolumne"}
VERDACHT = [
    (re.compile(r"\bWIP\b|\bTODO\b|\bTBD\b|\bDRAFT\b", re.I), "Arbeitsnotiz"),
    (re.compile(r"https?://"), "Quell-URL auf der Folie"),
    (re.compile(r"\bXXX+\b|\bplatzhalter\b|‹[^›]*›", re.I), "unersetzter Platzhalter"),
    (re.compile(r"^\s*(Lorem|Text hier)", re.I), "Blindtext"),
]


def cm(v):
    return round(Emu(v).cm, 1) if v is not None else None


def alle_shapes(shapes, tiefe=0):
    for sh in shapes:
        yield sh, tiefe
        if sh.shape_type == 6:  # GROUP
            yield from alle_shapes(sh.shapes, tiefe + 1)


def text_von(sh):
    return sh.text_frame.text.strip() if sh.has_text_frame else ""


def beschreibe(sh, tiefe, kurz=True):
    ind = "    " * tiefe
    if sh.is_placeholder:
        idx = sh.placeholder_format.idx
        rolle = ROLLEN.get(idx, "Inhaltsplatzhalter")
        adr = f"ph:{idx}"
        kopf = f"{ind}  {adr:<10} {rolle:<24}"
    else:
        adr = f"name:{sh.name}"
        kopf = f"{ind}  {adr:<10} {'':<24}" if kurz else f"{ind}  {adr}"
        kopf = f"{ind}  {'':<10} {sh.name[:24]:<24}"
    t = text_von(sh)
    if t:
        einzeilig = t.replace("\n", " | ")
        grenze = 78 if not kurz else 60
        t = einzeilig if len(einzeilig) <= grenze else einzeilig[:grenze] + "…"
        return f'{kopf} "{t}"'
    if sh.has_table:
        return f"{kopf} [Tabelle {len(sh.table.rows)}x{len(sh.table.columns)}]"
    if sh.has_chart:
        return f"{kopf} [Diagramm {sh.chart.chart_type}]"
    if sh.shape_type == 7 and "think-cell" in (sh.name or ""):
        return f"{kopf} [think-cell-Datenobjekt — nicht loeschen]"
    return None


def uebersicht(prs):
    print(f"Folien: {len(prs.slides)}   Groesse: {cm(prs.slide_width)} x {cm(prs.slide_height)} cm\n")
    for i, s in enumerate(prs.slides, 1):
        titel, kolumne = "", ""
        n_ph = n_leer = n_box = n_tab = n_chart = 0
        thinkcell = False
        for sh, _ in alle_shapes(s.shapes):
            if sh.is_placeholder:
                n_ph += 1
                idx = sh.placeholder_format.idx
                t = text_von(sh)
                if not t:
                    n_leer += 1
                if idx == 0:
                    titel = t
                elif idx == 13:
                    kolumne = t
            elif sh.has_text_frame and text_von(sh):
                n_box += 1
            if sh.has_table:
                n_tab += 1
            if sh.has_chart:
                n_chart += 1
            if "think-cell" in (sh.name or ""):
                thinkcell = True
        marker = " [think-cell]" if thinkcell else ""
        extra = f" Tab={n_tab}" if n_tab else ""
        extra += f" Chart={n_chart}" if n_chart else ""
        print(f"F{i:>2} | {s.slide_layout.name:<18} | PH {n_ph} (leer {n_leer}) Boxen {n_box}{extra}{marker}")
        titel_z = titel.replace("\n", " · ")
        print(f"     Titel:   {titel_z[:96] or '— kein Titel —'}")
        if kolumne:
            print(f"     Kolumne: {kolumne[:96]}")


def detail(prs, nr):
    if not 1 <= nr <= len(prs.slides):
        sys.exit(f"Folie {nr} gibt es nicht (1..{len(prs.slides)})")
    s = prs.slides[nr - 1]
    print(f"### Folie {nr} — Layout {s.slide_layout.name}\n")
    print("Platzhalter des Layouts:")
    for ph in s.slide_layout.placeholders:
        idx = ph.placeholder_format.idx
        print(f"  ph:{idx:<3} {ROLLEN.get(idx, 'Inhaltsplatzhalter')}")
    print("\nShapes der Folie:")
    for sh, tiefe in alle_shapes(s.shapes):
        zeile = beschreibe(sh, tiefe, kurz=False)
        if zeile:
            print(zeile)
    print("\nLeere Platzhalter:")
    leer = [
        f"ph:{sh.placeholder_format.idx} ({ROLLEN.get(sh.placeholder_format.idx, 'Inhalt')})"
        for sh, _ in alle_shapes(s.shapes)
        if sh.is_placeholder and not text_von(sh)
    ]
    print("  " + (", ".join(leer) if leer else "keine"))


def check(prs):
    befunde = []
    for i, s in enumerate(prs.slides, 1):
        layout = s.slide_layout.name
        ist_inhalt = "Inhalt" in layout
        vorhanden = {}
        for sh, _ in alle_shapes(s.shapes):
            if sh.is_placeholder:
                vorhanden[sh.placeholder_format.idx] = text_von(sh)
            t = text_von(sh)
            if t:
                for muster, was in VERDACHT:
                    if muster.search(t):
                        auszug = t.replace("\n", " ")[:70]
                        befunde.append((i, was, auszug))
                        break
        if ist_inhalt:
            for idx, name in PFLICHT_INHALT.items():
                if not vorhanden.get(idx, "").strip():
                    befunde.append((i, f"{name} leer (ph:{idx})", ""))
            # Die Quellenzeile ist Pflicht, wo Daten gezeigt werden oder die Folie
            # zu einem Inhaltskapitel gehoert. Credentials-Folien brauchen sie nicht.
            zeigt_daten = any(
                sh.has_table or sh.has_chart for sh, _ in alle_shapes(s.shapes)
            )
            kolumne = vorhanden.get(13, "").strip()
            braucht_quelle = zeigt_daten or (kolumne and not kolumne.startswith("1."))
            if braucht_quelle and not vorhanden.get(14, "").strip():
                befunde.append((i, "Quellenzeile leer (ph:14)", ""))
            titel = vorhanden.get(0, "")
            if titel and len(titel.split()) <= 3:
                befunde.append((i, "Titel ist ein Etikett, kein Aussagesatz", titel))

    kolumnen = {}
    for i, s in enumerate(prs.slides, 1):
        for sh, _ in alle_shapes(s.shapes):
            if sh.is_placeholder and sh.placeholder_format.idx == 13:
                t = text_von(sh)
                if t:
                    kolumnen.setdefault(t, []).append(i)
    if kolumnen:
        print("Kapitelkolumnen (muessen wortgleich zur Agenda sein):")
        for t, folien in sorted(kolumnen.items(), key=lambda x: x[1][0]):
            print(f"  F{','.join(map(str, folien)):<14} {t}")
        print()

    if not befunde:
        print("Pruefung ohne Befund.")
        return 0
    print(f"{len(befunde)} Befund(e):")
    for nr, was, auszug in befunde:
        print(f"  F{nr:>2}  {was}" + (f'  →  "{auszug}"' if auszug else ""))
    return 1


def main():
    p = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    p.add_argument("pptx")
    p.add_argument("--slide", type=int, help="eine Folie im Detail")
    p.add_argument("--check", action="store_true", help="Pruefmodus vor Uebergabe")
    a = p.parse_args()
    prs = Presentation(a.pptx)
    if a.slide:
        detail(prs, a.slide)
    elif a.check:
        sys.exit(check(prs))
    else:
        uebersicht(prs)


if __name__ == "__main__":
    main()
