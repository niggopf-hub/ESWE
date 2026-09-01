#!/usr/bin/env python3
"""Schreibt Text in VORHANDENE Boxen einer Unterlage — und legt niemals neue an.

Der Sinn: Platzhalter erben Schrift, Groesse, Farbe und Position vom Master. Ein neu
eingefuegtes Textfeld erbt nichts davon und macht die Datei unbrauchbar fuer die
Weiterbearbeitung. Deshalb bricht dieses Skript ab, wenn ein Ziel nicht existiert,
statt es anzulegen — ein Tippfehler soll auffallen, nicht stillschweigend eine Box
erzeugen.

    python3 fill_deck.py deck.pptx --map befuellung.json --out deck_v2.pptx
    python3 fill_deck.py deck.pptx --map befuellung.json --dry-run

Format der Zuordnungsdatei (Folienummer 1-basiert):

    {
      "12": {
        "ph:0":  "Starke Marktposition mit umfangreichem Investitionsprogramm …",
        "ph:13": "2. Musterversorger AG – Uebersicht und Herausforderungen",
        "ph:14": "Quellen: eigene Recherche, Jahresabschluss Musterversorger AG 2025",
        "ph:15": ["1) Net Debt = …", "2) Stichtag 31.12.2025"]
      },
      "13": { "name:Text Box 29": "999 Mio. EUR" }
    }

Eine Liste wird zu mehreren Absaetzen. Die Formatierung des ersten vorhandenen
Absatzes wird auf alle neuen Absaetze uebertragen, damit Aufzaehlungszeichen und
Schriftgrad erhalten bleiben.

Benoetigt: pip install python-pptx
"""
import argparse
import copy
import json
import sys

try:
    from pptx import Presentation
except ImportError:
    sys.exit("Fehlt: pip install python-pptx")


def alle_shapes(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:  # GROUP
            yield from alle_shapes(sh.shapes)


def finde(slide, adresse):
    """Loest 'ph:<idx>' oder 'name:<Shapename>' auf. Gibt None zurueck, wenn es fehlt."""
    art, _, wert = adresse.partition(":")
    if art == "ph":
        try:
            idx = int(wert)
        except ValueError:
            return None
        for sh in alle_shapes(slide.shapes):
            if sh.is_placeholder and sh.placeholder_format.idx == idx:
                return sh
    elif art == "name":
        for sh in alle_shapes(slide.shapes):
            if sh.name == wert:
                return sh
    return None


def schreibe(shape, inhalt):
    """Ersetzt den Text und erhaelt die Absatzformatierung des Originals."""
    tf = shape.text_frame
    zeilen = inhalt if isinstance(inhalt, list) else [inhalt]
    zeilen = [str(z) for z in zeilen]

    vorlage = copy.deepcopy(tf.paragraphs[0]._p)

    # Alle Absaetze bis auf den ersten entfernen, dann den ersten neu befuellen.
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)

    erster = tf.paragraphs[0]
    for r in list(erster.runs)[1:]:
        r._r.getparent().remove(r._r)
    if erster.runs:
        erster.runs[0].text = zeilen[0]
    else:
        erster.add_run().text = zeilen[0]

    for zeile in zeilen[1:]:
        neu = copy.deepcopy(vorlage)
        tf._txBody.append(neu)
        p = tf.paragraphs[-1]
        for r in list(p.runs)[1:]:
            r._r.getparent().remove(r._r)
        if p.runs:
            p.runs[0].text = zeile
        else:
            p.add_run().text = zeile


def main():
    ap = argparse.ArgumentParser(
        description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter
    )
    ap.add_argument("pptx")
    ap.add_argument("--map", required=True, help="JSON mit der Zuordnung")
    ap.add_argument("--out", help="Zieldatei (ohne --out nur Probelauf)")
    ap.add_argument("--dry-run", action="store_true", help="nur pruefen, nichts schreiben")
    a = ap.parse_args()

    prs = Presentation(a.pptx)
    with open(a.map, encoding="utf-8") as f:
        zuordnung = json.load(f)

    fehler, geschrieben = [], []
    for folie_nr, felder in zuordnung.items():
        try:
            nr = int(folie_nr)
        except ValueError:
            fehler.append(f"Folienummer '{folie_nr}' ist keine Zahl")
            continue
        if not 1 <= nr <= len(prs.slides):
            fehler.append(f"Folie {nr} gibt es nicht (1..{len(prs.slides)})")
            continue
        slide = prs.slides[nr - 1]
        for adresse, inhalt in felder.items():
            shape = finde(slide, adresse)
            if shape is None:
                fehler.append(
                    f"F{nr}: '{adresse}' nicht gefunden — "
                    "keine neue Box angelegt. Ziel mit inspect_deck.py pruefen."
                )
                continue
            if not shape.has_text_frame:
                fehler.append(f"F{nr}: '{adresse}' nimmt keinen Text auf")
                continue
            if not (a.dry_run or not a.out):
                schreibe(shape, inhalt)
            vorschau = (inhalt if isinstance(inhalt, str) else " | ".join(map(str, inhalt)))[:58]
            geschrieben.append(f"F{nr:>2}  {adresse:<22} ← {vorschau}")

    for z in geschrieben:
        print(z)
    if fehler:
        print(f"\n{len(fehler)} Problem(e):", file=sys.stderr)
        for f in fehler:
            print(f"  {f}", file=sys.stderr)
        sys.exit(1)

    if a.out and not a.dry_run:
        prs.save(a.out)
        print(f"\nGespeichert: {a.out}  ({len(geschrieben)} Felder)")
    else:
        print(f"\nProbelauf ohne Befund ({len(geschrieben)} Felder wuerden geschrieben).")


if __name__ == "__main__":
    main()
