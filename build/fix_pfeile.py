# -*- coding: utf-8 -*-
"""Verbindungslinien auf der Beteiligungsstruktur und dem One Pager begradigen.

Beim Einsetzen der Logos sind mehrere Konnektoren umgeklappt (waagerecht statt
senkrecht) bzw. mit kleinem Versatz stehen geblieben. Hier werden sie auf exakt
senkrechte bzw. waagerechte Strecken gesetzt - Breite oder Hoehe genau null,
kein Knickwert. Eine Strecke ohne Versatz zeichnet jeder Renderer gleich.
"""
import copy, sys
from pptx import Presentation
from pptx.oxml.ns import qn

EMU = 12700.0
SRC = sys.argv[1] if len(sys.argv) > 1 else 'user_v1.pptx'
DST = sys.argv[2] if len(sys.argv) > 2 else 'user_v1_fix.pptx'


def cxns(slide):
    out = []
    def walk(shapes):
        for sh in shapes:
            if sh.shape_type == 6:
                walk(sh.shapes)
            elif sh._element.tag.endswith('}cxnSp'):
                out.append(sh)
    walk(slide.shapes)
    return out


def geom(sh):
    x = sh._element.spPr.find(qn('a:xfrm'))
    o, e = x.find(qn('a:off')), x.find(qn('a:ext'))
    return (int(o.get('x')) / EMU, int(o.get('y')) / EMU,
            int(e.get('cx')) / EMU, int(e.get('cy')) / EMU)


def set_line(sh, x0, y0, x1, y1):
    """Strecke exakt achsparallel setzen und Knickwert entfernen."""
    x = sh._element.spPr.find(qn('a:xfrm'))
    for a in ('flipH', 'flipV', 'rot'):
        if x.get(a) is not None:
            del x.attrib[a]
    if x1 < x0:
        x0, x1 = x1, x0
        x.set('flipH', '1')
    if y1 < y0:
        y0, y1 = y1, y0
        x.set('flipV', '1')
    o, e = x.find(qn('a:off')), x.find(qn('a:ext'))
    o.set('x', str(int(round(x0 * EMU))))
    o.set('y', str(int(round(y0 * EMU))))
    e.set('cx', str(int(round((x1 - x0) * EMU))))
    e.set('cy', str(int(round((y1 - y0) * EMU))))
    av = sh._element.spPr.find(qn('a:prstGeom')).find(qn('a:avLst'))
    for c in list(av):
        av.remove(c)


def clone_after(sh):
    """Zweites Segment mit identischer Linienformatierung anlegen."""
    new = copy.deepcopy(sh._element)
    nv = new.find(qn('p:nvCxnSpPr')).find(qn('p:cNvCxnSpPr'))
    for tag in ('a:stCxn', 'a:endCxn'):
        e = nv.find(qn(tag))
        if e is not None:
            nv.remove(e)
    cNvPr = new.find(qn('p:nvCxnSpPr')).find(qn('p:cNvPr'))
    cNvPr.set('id', str(9000 + clone_after.n))
    cNvPr.set('name', 'Verbinder %d' % (9000 + clone_after.n))
    clone_after.n += 1
    sh._element.addnext(new)
    from pptx.shapes.connector import Connector
    return Connector(new, sh._parent)
clone_after.n = 0


def unglue(sh):
    nv = sh._element.find(qn('p:nvCxnSpPr')).find(qn('p:cNvCxnSpPr'))
    for tag in ('a:stCxn', 'a:endCxn'):
        e = nv.find(qn(tag))
        if e is not None:
            nv.remove(e)


prs = Presentation(SRC)

# ---------------------------------------------------------------- Folie 14
# Sollwerte: senkrechte Abgaenge auf den Spaltenmitten, senkrechte Klammer
# in der Gasse zwischen Spalte 3 und 4.
SOLL14 = [
    # (Erkennung: aktuelle x/y ungefaehr) -> (x, y_oben, y_unten)
    ((416.0, 147.0), (421.0, 142.0, 152.0)),    # Stadt -> WVV
    ((414.0, 183.0), (421.0, 176.0, 190.0)),    # WVV -> ESWE
    ((765.0, 183.0), (772.5, 176.0, 190.0)),    # Thuega -> ESWE
    ((79.1, 242.9), (90.2, 232.0, 254.0)),      # Abgaenge Reihe 1
    ((212.0, 243.0), (222.5, 232.0, 254.0)),
    ((344.4, 243.0), (354.9, 232.0, 254.0)),
    ((609.0, 243.0), (619.5, 232.0, 254.0)),
    ((348.9, 306.0), (421.0, 232.0, 380.0)),    # Klammer zur zweiten Reihe
    ((82.7, 388.0), (90.2, 380.0, 396.0)),      # Abgang Reihe 2, Spalte 1
]

s14 = prs.slides[13]
fixed = 0
for sh in cxns(s14):
    x, y, w, h = geom(sh)
    for (rx, ry), (nx, y0, y1) in SOLL14:
        if abs(x - rx) < 2.0 and abs(y - ry) < 2.0:
            set_line(sh, nx, y0, nx, y1)
            fixed += 1
            break
print('Folie 14: %d Verbindungslinien begradigt' % fixed)

# ---------------------------------------------------------------- Folie 13
# Donut: Mittelpunkt und sichtbarer Aussenradius; Callouts docken achsparallel an.
CX, CY, R = 562.4, 196.7, 66.0
import math


def rand_x(y, links=True):
    dy = abs(y - CY)
    if dy >= R:
        return CX
    d = math.sqrt(R * R - dy * dy)
    return CX - d if links else CX + d


def rand_y(x, oben=False):
    dx = abs(x - CX)
    if dx >= R:
        return CY
    d = math.sqrt(R * R - dx * dx)
    return CY - d if oben else CY + d


s13 = prs.slides[12]
cs = {round(geom(c)[0], 1): c for c in cxns(s13)}
todo = 0
for key, c in list(cs.items()):
    x, y, w, h = geom(c)
    if w < 0.6 or h < 0.6:      # die senkrechten Trenner der Icon-Leiste
        continue
    unglue(c)
    if abs(x - 462.8) < 2:                       # Stromverkauf -> Donut, waagerecht
        yy = 214.8
        set_line(c, 462.7, yy, rand_x(yy, True), yy)
    elif abs(x - 621.1) < 2:                     # Wasser -> Donut, zwei Segmente
        set_line(c, 732.3, 282.6, 732.3, 258.0)
        c2 = clone_after(c)
        set_line(c2, 732.3, 258.0, rand_x(258.0, False), 258.0)
    elif abs(x - 638.8) < 2:                     # Dienstleistungen -> Donut, waagerecht
        yy = 214.8
        set_line(c, 661.4, yy, rand_x(yy, False), yy)
    elif abs(x - 391.9) < 2:                     # Gasverkauf -> Donut, zwei Segmente
        set_line(c, 391.9, 282.6, 391.9, 258.0)
        c2 = clone_after(c)
        set_line(c2, 391.9, 258.0, rand_x(258.0, True), 258.0)
    elif abs(x - 551.5) < 2:                     # Waerme -> Donut, senkrecht
        set_line(c, 562.1, 282.6, 562.1, rand_y(562.1, False))
    else:
        continue
    todo += 1
print('Folie 13: %d Anschlusslinien neu gefuehrt' % todo)

prs.save(DST)
print('gespeichert:', DST)
