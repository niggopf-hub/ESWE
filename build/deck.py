# -*- coding: utf-8 -*-
"""Bau-Framework: repliziert das evm-Layout (A4 quer, 842 x 595,32 pt) in PPTX."""
import os, json
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

BASE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(BASE, 'assets')

PW, PH = 842.04, 595.32                      # Seitenmaße in pt (A4 quer)

# --- Metzler Farbwelt (aus dem evm-Deck extrahiert) ---
DARK   = RGBColor(0x00, 0x3D, 0x7C)          # Metzler-Dunkelblau
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x80, 0x80, 0x80)
LIGHTBOX = RGBColor(0xEB, 0xF4, 0xFF)          # heller Kastenhintergrund
LIGHT  = RGBColor(0xD9, 0xE4, 0xEF)
MID    = RGBColor(0x5B, 0x8D, 0xB8)          # mittleres Blau (Charts)
ORANGE = RGBColor(0xE8, 0x7B, 0x1E)

# --- Schriften (Originalnamen des Metzler-Templates) ---
F_LIGHT = 'Univers for Metzler Light'
F_REG   = 'Univers for Metzler'
F_TWO   = 'Univers for Metzler 2'
F_WING  = 'Wingdings'

# QA-Modus: Wingdings-Bullets durch darstellbares Zeichen ersetzen
QA = os.environ.get('DECK_QA') == '1'

FONTMAP = {
    'UniversForMetzler-Light':        (F_LIGHT, False, False),
    'UniversForMetzler-LightItalic':  (F_LIGHT, False, True),
    'UniversForMetzler-Regula':       (F_REG,   False, False),
    'UniversForMetzler-Regular':      (F_REG,   False, False),
    'UniversForMetzler-Bold':         (F_REG,   True,  False),
    'UniversForMetzler-BoldIt':       (F_REG,   True,  True),
    'UniversForMetzler-BoldItalic':   (F_REG,   True,  True),
    'UniversForMetzlerTwo-Bol':       (F_TWO,   True,  False),
    'UniversForMetzlerTwo-Bold':      (F_TWO,   True,  False),
    'UniversForMetzlerTwo-BoldIt':    (F_TWO,   True,  True),
    'UniversForMetzlerTwo-BoldItalic':(F_TWO,   True,  True),
    'Wingdings-Regular':              (F_WING,  False, False),
    'Wingdings3':                     ('Wingdings 3', False, False),
    'ArialMT':                        ('Arial',  False, False),
    'CambriaMath':                    ('Cambria Math', False, False),
}

BULLET = 'n'          # Wingdings 'n' == ◼ (Metzler-Aufzählungszeichen)
DASH   = '–'


def rgb(v):
    return RGBColor((v >> 16) & 255, (v >> 8) & 255, v & 255)


def new_deck():
    prs = Presentation()
    prs.slide_width = Pt(PW)
    prs.slide_height = Pt(PH)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------- Textbausteine
class Run(object):
    """Ein Textlauf mit eigener Formatierung."""
    def __init__(self, text, font=F_LIGHT, size=10.6, color=BLACK,
                 bold=False, italic=False, baseline=None):
        self.text, self.font, self.size = text, font, size
        self.color, self.bold, self.italic = color, bold, italic
        self.baseline = baseline          # z. B. 30 fuer Hochstellung


def wing(size=10.6, color=DARK):
    """Metzler-Bullet (Wingdings ◼)."""
    if QA:
        return Run('■', 'Liberation Sans', size * 0.85, color)
    return Run(BULLET, F_WING, size, color)


def _apply(run, r):
    run.text = r.text
    f = run.font
    f.name = r.font
    f.size = Pt(r.size)
    f.bold = r.bold
    f.italic = r.italic
    f.color.rgb = r.color
    if getattr(r, 'baseline', None):
        run._r.get_or_add_rPr().set('baseline', str(int(r.baseline * 1000)))
    # Ost-/Komplexschrift gleich setzen, damit PowerPoint nicht substituiert
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set('typeface', r.font)


def _set_spacing(p, line=None, space_before=None, space_after=None):
    pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
    if line is not None:
        ln = pPr.find(qn('a:lnSpc'))
        if ln is None:
            ln = pPr.makeelement(qn('a:lnSpc'), {})
            pPr.insert(0, ln)
        for c in list(ln):
            ln.remove(c)
        pts = ln.makeelement(qn('a:spcPts'), {})
        pts.set('val', str(int(round(line * 100))))
        ln.append(pts)
    if space_before is not None:
        p.space_before = Pt(space_before)
    if space_after is not None:
        p.space_after = Pt(space_after)


def textbox(slide, x, y, w, h, paras, anchor='t', align=PP_ALIGN.LEFT,
            wrap=True, line=None, space_before=0, space_after=0,
            indent=None, hanging=None):
    """paras: Liste von Absätzen; jeder Absatz ist ein Run, ein String oder eine Run-Liste.
    Ein Absatz darf auch ein dict sein: {'runs':[...], 'line':..,'before':..,'indent':..,'hanging':..,'align':..}"""
    box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {'t': MSO_ANCHOR.TOP, 'm': MSO_ANCHOR.MIDDLE,
                          'b': MSO_ANCHOR.BOTTOM}[anchor]
    box.fill.background()
    box.line.fill.background()
    first = True
    for para in paras:
        opts = {}
        if isinstance(para, dict):
            opts = para
            runs = opts.get('runs', [])
        else:
            runs = para
        if isinstance(runs, (Run, str)):
            runs = [runs]
        runs = [Run(r) if isinstance(r, str) else r for r in runs]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opts.get('align', align)
        _set_spacing(p, opts.get('line', line),
                     opts.get('before', space_before), opts.get('after', space_after))
        ind = opts.get('indent', indent)
        hang = opts.get('hanging', hanging)
        if ind is not None or hang is not None:
            pPr = p._p.get_or_add_pPr()
            if ind is not None:
                pPr.set('marL', str(int(round(ind * 12700))))
            if hang is not None:
                pPr.set('indent', str(int(round(-hang * 12700))))
        for r in runs:
            _apply(p.add_run(), r)
    return box


def line_text(slide, x, y0, y1, text, font=F_LIGHT, size=10.6, color=BLACK,
              bold=False, italic=False, align=PP_ALIGN.LEFT, w=None):
    """Einzelne Textzeile exakt auf der Höhe der Original-Zeile (mittig verankert)."""
    h = (y1 - y0)
    w = w if w is not None else 700
    return textbox(slide, x, y0 - 3, w, h + 6,
                   [[Run(text, font, size, color, bold, italic)]],
                   anchor='m', align=align, wrap=False)


# ---------------------------------------------------------------- Grafik
def rect(slide, x, y, w, h, fill=None, linecolor=None, linew=0.75, shape=MSO_SHAPE.RECTANGLE,
         adj=None):
    s = slide.shapes.add_shape(shape, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if linecolor is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = linecolor
        s.line.width = Pt(linew)
    no_shadow(s)
    strip_style(s)
    if adj is not None:
        try:
            s.adjustments[0] = adj
        except Exception:
            pass
    s.text_frame.word_wrap = False
    return s


def hline(slide, x0, x1, y, color=DARK, width=0.75, dash=None):
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    c = slide.shapes.add_connector(1, Pt(x0), Pt(y), Pt(x1), Pt(y))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dash:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return c


def vline(slide, x, y0, y1, color=DARK, width=0.75, dash=None):
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    c = slide.shapes.add_connector(1, Pt(x), Pt(y0), Pt(x), Pt(y1))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dash:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return c


def picture(slide, name, x, y, w, h):
    return slide.shapes.add_picture(os.path.join(ASSETS, name), Pt(x), Pt(y), Pt(w), Pt(h))


# ---------------------------------------------------------------- Seitengerüst
LOGO = 'img_6.png'          # Metzler-Logo unten rechts


def chrome(slide, number, header=None, sources=None, footnotes=None,
           rule_y=92.0, logo=True):
    """Kopfzeile, Trennlinie, Fußnoten, Quellenzeile, Seitenzahl, Logo."""
    if rule_y:
        hline(slide, 30, 812, rule_y)
    if header:
        line_text(slide, 29.8, 20.9, 30.9, header, F_REG, 10, GREY)
    y = 563.5
    if footnotes:
        textbox(slide, 29.8, 543.0, 760, 20,
                [[Run(footnotes, F_LIGHT, 6.5, BLACK)]], anchor='b',
                line=7.5, wrap=True)
    if sources:
        line_text(slide, 47.5, 563.7, 571.8, sources, F_LIGHT, 8.1, BLACK, w=640)
    line_text(slide, 30.0, 563.5, 571.6, str(number), F_LIGHT, 8.0, BLACK, w=30)
    # Fußlinie
    hline(slide, 30, 689, 557)
    hline(slide, 699, 786, 557)
    hline(slide, 796, 812, 557)
    if logo:
        picture(slide, LOGO, 699, 543, 111, 28)


def title_2lines(slide, l1, l2=None, size=20.0, y=40.4, color=BLACK, font=None):
    """Slide-Headline im evm-Muster: 20 pt fett, Zeilenraster 24 pt."""
    font = font or 'UniversForMetzler-Bold'
    fam, bold, ital = FONTMAP[font]
    line_text(slide, 29.8, y, y + 20.1, l1, fam, size, color, bold, ital, w=790)
    if l2:
        line_text(slide, 29.8, y + 24.0, y + 44.1, l2, fam, size, color, bold, ital, w=790)


# ---------------------------------------------------------------- Textmaß
from PIL import ImageFont
_FDIR = '/usr/share/fonts/truetype/liberation/'
_CACHE = {}


def _metric_font(size, bold=False, italic=False):
    key = (round(size, 2), bold, italic)
    if key not in _CACHE:
        sfx = 'Bold' if bold and not italic else ('BoldItalic' if bold and italic
              else ('Italic' if italic else 'Regular'))
        _CACHE[key] = ImageFont.truetype(_FDIR + 'LiberationSans-%s.ttf' % sfx,
                                         int(round(size * 8)))
    return _CACHE[key]


def tw(s, size, bold=False, italic=False):
    """Textbreite in pt (Liberation-Sans-Metrik als Univers-Näherung)."""
    return _metric_font(size, bold, italic).getlength(s) / 8.0


def wrap(text, width, size, bold=False, italic=False):
    """Greedy-Umbruch auf feste Breite; liefert Liste von Zeilen."""
    out, cur = [], ''
    for word in text.split(' '):
        cand = word if not cur else cur + ' ' + word
        if tw(cand, size, bold, italic) <= width or not cur:
            cur = cand
        else:
            out.append(cur)
            cur = word
    if cur:
        out.append(cur)
    return out


def para_lines(text, width, size, bold=False, italic=False):
    """Wie wrap(), aber respektiert manuelle Umbrüche mit '|'."""
    res = []
    for chunk in text.split('|'):
        res.extend(wrap(chunk.strip(), width, size, bold, italic))
    return res



def set_alpha(shape, pct):
    """Fuellung teiltransparent setzen (pct = Deckkraft in Prozent)."""
    from pptx.oxml import parse_xml
    solidFill = shape._element.spPr.find(qn('a:solidFill'))
    clr = solidFill.find(qn('a:srgbClr'))
    for c in list(clr):
        clr.remove(c)
    clr.append(parse_xml('<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml'
                         '/2006/main" val="%d"/>' % int(pct * 1000)))
    return shape


def no_shadow(shape):
    from pptx.oxml import parse_xml
    spPr = shape._element.spPr
    old = spPr.find(qn('a:effectLst'))
    if old is not None:
        spPr.remove(old)
    spPr.append(parse_xml('<a:effectLst xmlns:a="http://schemas.openxmlformats.org'
                          '/drawingml/2006/main"/>'))
    return shape


def soft_shadow(shape, blur=5.0, dist=2.0, dir_deg=45, alpha=22):
    """Dezenter Schlagschatten wie in der Metzler-Vorlage."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsmap
    spPr = shape._element.spPr
    xml = (
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:outerShdw blurRad="%d" dist="%d" dir="%d" rotWithShape="0">'
        '<a:srgbClr val="000000"><a:alpha val="%d"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'
        % (int(blur * 12700), int(dist * 12700), int(dir_deg * 60000), int(alpha * 1000))
    )
    for tag in ('a:effectLst',):
        old = spPr.find(qn(tag))
        if old is not None:
            spPr.remove(old)
    spPr.append(parse_xml(xml))
    return shape


NB = ' '


def nb(s):
    """Schützt stehende Wendungen gegen Zeilenumbruch."""
    for a, b in (('p. a.', 'p.' + NB + 'a.'), ('z. B.', 'z.' + NB + 'B.'),
                 ('u. a.', 'u.' + NB + 'a.'), ('i. H. v.', 'i.' + NB + 'H.' + NB + 'v.'),
                 ('Mio. EUR', 'Mio.' + NB + 'EUR'), ('ca. ', 'ca.' + NB),
                 ('Ø ~', 'Ø' + NB + '~')):
        s = s.replace(a, b)
    return s


# ---------------------------------------------------------------- Donut
def block_arc(slide, cx, cy, r_out, r_in, start_deg, end_deg, fill, steps=None):
    """Donut-Segment als Freiform-Polygon (renderer-unabhaengig exakt)."""
    import math
    sweep = end_deg - start_deg
    steps = steps or max(6, int(abs(sweep) / 2.0) + 2)
    pts = []
    for i in range(steps + 1):
        a = math.radians(start_deg + sweep * i / float(steps))
        pts.append((cx + r_out * math.cos(a), cy + r_out * math.sin(a)))
    for i in range(steps + 1):
        a = math.radians(end_deg - sweep * i / float(steps))
        pts.append((cx + r_in * math.cos(a), cy + r_in * math.sin(a)))
    b = slide.shapes.build_freeform(Pt(pts[0][0]), Pt(pts[0][1]))
    b.add_line_segments([(Pt(x), Pt(y)) for x, y in pts[1:]], close=True)
    sh = b.convert_to_shape()
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    no_shadow(sh)
    return sh


def _block_arc_preset(slide, cx, cy, r_out, r_in, start_deg, end_deg, fill):
    """(ungenutzt) Preset-Variante""" 
    s = slide.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, Pt(cx - r_out), Pt(cy - r_out),
                               Pt(2 * r_out), Pt(2 * r_out))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    prstGeom = s._element.spPr.find(qn('a:prstGeom'))
    av = prstGeom.find(qn('a:avLst'))
    for c in list(av):
        av.remove(c)
    for name, val in (('adj1', int(round(start_deg * 60000))),
                      ('adj2', int(round(end_deg * 60000))),
                      ('adj3', int(round((1 - r_in / float(r_out)) / 2.0 * 100000)))):
        gd = av.makeelement(qn('a:gd'), {'name': name, 'fmla': 'val %d' % val})
        av.append(gd)
    s.text_frame.word_wrap = False
    return s


def oval(slide, cx, cy, w, h, fill=None, linecolor=None, linew=0.75):
    return rect(slide, cx - w / 2.0, cy - h / 2.0, w, h, fill=fill,
                linecolor=linecolor, linew=linew, shape=MSO_SHAPE.OVAL)


def dashed_rect(slide, x, y, w, h, linecolor, fill=None, linew=0.75):
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    s = rect(slide, x, y, w, h, fill=fill, linecolor=linecolor, linew=linew)
    s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return s


def elbow(slide, pts, color, width=0.75):
    """Polylinie aus Punktliste (Verbindungslinien wie in der Vorlage)."""
    out = []
    for i in range(len(pts) - 1):
        (x0, y0), (x1, y1) = pts[i], pts[i + 1]
        c = slide.shapes.add_connector(1, Pt(x0), Pt(y0), Pt(x1), Pt(y1))
        c.line.color.rgb = color
        c.line.width = Pt(width)
        out.append(c)
    return out


def vlabel(slide, x, y, w, h, paras, size=9.0, color=DARK, font=None, line=11.0,
           bold=True, rot=270):
    """Vertikal gesetzte Label-Kachel-Beschriftung (wie evm-Seitenkacheln)."""
    font = font or F_TWO
    cx, cy = x + w / 2.0, y + h / 2.0
    bw, bh = h, w
    box = textbox(slide, cx - bw / 2.0, cy - bh / 2.0, bw, bh,
                  [p if isinstance(p, list) else [Run(p, font, size, color, bold)]
                   for p in paras],
                  anchor='m', align=PP_ALIGN.CENTER, line=line, wrap=False)
    box.rotation = rot
    return box


# ---------------------------------------------------------------- Konnektoren
def strip_style(shape):
    """Theme-Style entfernen, damit Fuellung/Linie/Effekt nur aus spPr kommen."""
    st = shape._element.find(qn('p:style'))
    if st is not None:
        shape._element.remove(st)
    return shape


def connect(slide, x0, y0, x1, y1, begin=None, end=None, color=None,
            width=0.9, dash=None):
    """Winkelverbinder (PowerPoint: Verbinder – Winkel), an Formen geklebt.

    begin/end: (Form, idx) zum Ankleben – idx 0 = oben mittig, 1 = links,
    2 = unten mittig, 3 = rechts. Ohne Angabe bleibt das Ende frei.

    Die Geometrie wird bewusst immer waagerecht oder senkrecht gesetzt: ein
    Winkelverbinder ohne Versatz zeichnet eine gerade Strecke, und zwar in jedem
    Renderer gleich. Ein erzwungener Knick (adj1) wuerde in PowerPoint die
    senkrechte Teilstrecke verschieben und quer durch die Kaesten laufen.
    """
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    c = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW,
                                   Pt(x0), Pt(y0), Pt(x1), Pt(y1))
    if begin is not None:
        c.begin_connect(begin[0], begin[1])
    if end is not None:
        c.end_connect(end[0], end[1])
    c.line.color.rgb = color or DARK
    c.line.width = Pt(width)
    if dash:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    strip_style(c)
    return c


# ---------------------------------------------------------------- Aufzählungen
def _bullet_pPr(p, color, size, indent=11.0, before=None, line=None):
    """Echte PowerPoint-Aufzählung (Wingdings-Quadrat) am Absatz setzen."""
    from pptx.oxml import parse_xml
    A = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(int(round(indent * 12700))))
    pPr.set('indent', str(int(round(-indent * 12700))))
    _set_spacing(p, line, before, None)
    for xml in ('<a:buClr %s><a:srgbClr val="%02X%02X%02X"/></a:buClr>'
                % (A, color[0], color[1], color[2]),
                '<a:buSzPct %s val="75000"/>' % A,
                '<a:buFont %s typeface="%s"/>' % (A, 'Arial' if QA else F_WING),
                '<a:buChar %s char="%s"/>' % (A, '&#9642;' if QA else 'n')):
        pPr.append(parse_xml(xml))
    return p



def add_runs(paragraph, runs):
    """Weitere Textlaeufe an einen bestehenden Absatz anhaengen."""
    for r in runs:
        _apply(paragraph.add_run(), r)
    return paragraph


# ---------------------------------------------------------------- Shape-Text
def shape_text(shape, paras, size=9.0, font=None, color=BLACK, bold=False,
               italic=False, align=PP_ALIGN.LEFT, anchor='m', line=None,
               before=0.0, inset=(0.0, 0.0, 0.0, 0.0), bullets=False,
               bullet_color=None, indent=11.0, wrap=True, bullet_from=0):
    """Schreibt direkt in den Textrahmen einer Form – kein zusaetzliches Textfeld.

    paras: je Absatz ein String, ein Run oder eine Run-Liste.
    inset: (links, oben, rechts, unten) in pt.
    bullets=True setzt echte Aufzaehlungszeichen (Wingdings-Quadrat);
    bullet_from bestimmt, ab welchem Absatz (z. B. 1 fuer eine Ueberschrift ohne
    Aufzaehlungszeichen).
    """
    font = font or F_LIGHT
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left, tf.margin_top, tf.margin_right, tf.margin_bottom = (
        Pt(inset[0]), Pt(inset[1]), Pt(inset[2]), Pt(inset[3]))
    tf.vertical_anchor = {'t': MSO_ANCHOR.TOP, 'm': MSO_ANCHOR.MIDDLE,
                          'b': MSO_ANCHOR.BOTTOM}[anchor]
    first = True
    for idx, para in enumerate(paras):
        if isinstance(para, list):
            runs = para
        elif isinstance(para, Run):
            runs = [para]
        else:
            runs = [Run(para, font, size, color, bold, italic)]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = align
        if bullets and idx >= bullet_from:
            _bullet_pPr(p, bullet_color or DARK, size, indent,
                        before=(0 if first else before), line=line)
        else:
            _set_spacing(p, line, (0 if first else before), None)
        for r in runs:
            _apply(p.add_run(), r)
        first = False
    return shape
